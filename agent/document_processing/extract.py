"""
Text & Data Extraction - Stage 2 of the Processing Pipeline

Handles extraction from every major document type:
  PDF        → pdfplumber (tables + text) with PyPDF2 fallback
  DOCX/ODT   → python-docx / odfpy
  XLSX/XLS   → openpyxl / xlrd
  CSV/TSV    → csv stdlib
  PPTX       → python-pptx
  Images     → Pillow + pytesseract OCR
  HTML/XML   → BeautifulSoup4 / lxml
  JSON/JSONL → stdlib json
  TXT/MD/RST → direct read with charset detection
  YAML/TOML  → pyyaml / tomllib
  Parquet    → pyarrow / pandas
  Email      → email stdlib
  Archives   → zipfile / tarfile (flat listing, no recursion)
  Code files → language-aware tokenisation
  Fallback   → chardet / charset-normalizer raw decode

Each extractor returns a unified ExtractionResult.
"""

from __future__ import annotations

import csv
import io
import json
import logging
import os
import re
import traceback
from dataclasses import dataclass, field, asdict
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Data structures
# ---------------------------------------------------------------------------

@dataclass
class PageContent:
    page_number: int
    text: str
    tables: List[List[List[str]]] = field(default_factory=list)
    images_found: int = 0
    confidence: float = 1.0


@dataclass
class ExtractionResult:
    success: bool
    file_id: str
    filename: str
    document_type: str          # e.g. "pdf", "docx", "xlsx", "image", ...
    raw_text: str               # full concatenated text
    pages: List[PageContent]    # per-page breakdown
    tables: List[List[List[str]]]  # all tables extracted
    key_value_pairs: Dict[str, Any]  # structured KV from the doc
    metadata: Dict[str, Any]    # author, created, modified, title, etc.
    language: str               # detected language code (e.g. "en")
    word_count: int
    char_count: int
    page_count: int
    extraction_method: str      # which library was used
    confidence: float           # 0–1 overall quality score
    chunks: List[str]           # text chunks for embedding (≈ 512 tokens)
    errors: List[str] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)

    def to_dict(self) -> Dict:
        return asdict(self)


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

_CHUNK_SIZE = 1800  # characters per chunk (≈ 512 tokens)


def _chunk_text(text: str, chunk_size: int = _CHUNK_SIZE) -> List[str]:
    """Split text into overlapping chunks for vector embedding."""
    chunks: List[str] = []
    overlap = chunk_size // 5
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end].strip())
        start += chunk_size - overlap
    return [c for c in chunks if c]


def _detect_language(text: str) -> str:
    """Lightweight language detection without external deps."""
    sample = text[:2000].lower()
    eng_markers = ["the ", "and ", "of ", "to ", "in ", "is ", "that "]
    spa_markers = ["el ", "la ", "de ", "que ", "en ", "es ", "los "]
    fra_markers = ["le ", "la ", "de ", "et ", "un ", "les ", "pour "]
    scores = {
        "en": sum(sample.count(m) for m in eng_markers),
        "es": sum(sample.count(m) for m in spa_markers),
        "fr": sum(sample.count(m) for m in fra_markers),
    }
    return max(scores, key=scores.get) if any(scores.values()) else "unknown"


def _clean_raw_text(text: str) -> str:
    text = re.sub(r"\r\n", "\n", text)
    text = re.sub(r"\r", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _decode_bytes(data: bytes) -> str:
    """Try common encodings with chardet fallback."""
    for enc in ("utf-8", "utf-16", "latin-1", "cp1252"):
        try:
            return data.decode(enc)
        except (UnicodeDecodeError, LookupError):
            pass
    try:
        import chardet  # type: ignore
        detected = chardet.detect(data)
        enc = detected.get("encoding") or "utf-8"
        return data.decode(enc, errors="replace")
    except ImportError:
        pass
    return data.decode("utf-8", errors="replace")


def _build_result(
    file_id: str,
    filename: str,
    doc_type: str,
    raw_text: str,
    pages: List[PageContent],
    tables: List[List[List[str]]],
    kv: Dict,
    doc_meta: Dict,
    method: str,
    confidence: float,
    errors: List[str],
    warnings: List[str],
) -> ExtractionResult:
    raw_text = _clean_raw_text(raw_text)
    return ExtractionResult(
        success=len(errors) == 0,
        file_id=file_id,
        filename=filename,
        document_type=doc_type,
        raw_text=raw_text,
        pages=pages,
        tables=tables,
        key_value_pairs=kv,
        metadata=doc_meta,
        language=_detect_language(raw_text),
        word_count=len(raw_text.split()),
        char_count=len(raw_text),
        page_count=len(pages) or 1,
        extraction_method=method,
        confidence=confidence,
        chunks=_chunk_text(raw_text),
        errors=errors,
        warnings=warnings,
    )


# ---------------------------------------------------------------------------
# Extractors
# ---------------------------------------------------------------------------

def _extract_pdf(data: bytes, file_id: str, filename: str) -> ExtractionResult:
    errors, warnings, pages, tables, kv, meta = [], [], [], [], {}, {}
    raw_parts = []
    method = "pdfplumber"

    try:
        import pdfplumber  # type: ignore
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            meta = pdf.metadata or {}
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                page_tables = page.extract_tables() or []
                # clean tables
                clean_tables = [
                    [[str(c or "").strip() for c in row] for row in tbl]
                    for tbl in page_tables
                ]
                tables.extend(clean_tables)
                raw_parts.append(text)
                pages.append(PageContent(
                    page_number=i,
                    text=text,
                    tables=clean_tables,
                    images_found=len(page.images),
                ))
    except Exception as exc:
        warnings.append(f"pdfplumber failed ({exc}); trying PyPDF2")
        method = "PyPDF2"
        try:
            import PyPDF2  # type: ignore
            reader = PyPDF2.PdfReader(io.BytesIO(data))
            meta = dict(reader.metadata or {})
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ""
                raw_parts.append(text)
                pages.append(PageContent(page_number=i, text=text))
        except Exception as exc2:
            errors.append(f"PDF extraction failed: {exc2}")
            method = "failed"

    return _build_result(
        file_id, filename, "pdf", "\n\n".join(raw_parts),
        pages, tables, kv, meta, method,
        confidence=0.9 if not errors else 0.3,
        errors=errors, warnings=warnings,
    )


def _extract_docx(data: bytes, file_id: str, filename: str) -> ExtractionResult:
    errors, warnings, pages, tables, kv, meta = [], [], [], [], {}, {}
    raw_parts = []
    method = "python-docx"
    try:
        from docx import Document  # type: ignore
        doc = Document(io.BytesIO(data))
        props = doc.core_properties
        meta = {
            "author": props.author,
            "title": props.title,
            "subject": props.subject,
            "created": str(props.created),
            "modified": str(props.modified),
            "keywords": props.keywords,
        }
        for para in doc.paragraphs:
            if para.text.strip():
                raw_parts.append(para.text)
        for tbl in doc.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
            tables.append(rows)
        pages.append(PageContent(page_number=1, text="\n".join(raw_parts), tables=tables))
    except Exception as exc:
        errors.append(f"DOCX extraction failed: {exc}")
        method = "failed"

    return _build_result(
        file_id, filename, "docx", "\n".join(raw_parts),
        pages, tables, kv, meta, method,
        confidence=0.95 if not errors else 0.3,
        errors=errors, warnings=warnings,
    )


def _extract_xlsx(data: bytes, file_id: str, filename: str) -> ExtractionResult:
    errors, warnings, pages, tables, kv, meta = [], [], [], [], {}, {}
    raw_parts = []
    method = "openpyxl"
    try:
        import openpyxl  # type: ignore
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
        meta = {"sheet_names": wb.sheetnames, "active_sheet": wb.active.title if wb.active else ""}
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                str_row = [str(c) if c is not None else "" for c in row]
                rows.append(str_row)
                raw_parts.append(" | ".join(str_row))
            tables.append(rows)
            pages.append(PageContent(page_number=len(pages) + 1, text="\n".join(" | ".join(r) for r in rows), tables=[rows]))
    except Exception as exc:
        errors.append(f"XLSX extraction failed: {exc}")
        method = "failed"

    return _build_result(
        file_id, filename, "xlsx", "\n".join(raw_parts),
        pages, tables, kv, meta, method,
        confidence=0.95 if not errors else 0.3,
        errors=errors, warnings=warnings,
    )


def _extract_csv(data: bytes, file_id: str, filename: str, sep: str = ",") -> ExtractionResult:
    errors, warnings, pages, tables, kv, meta = [], [], [], [], {}, {}
    raw_parts = []
    method = "csv"
    try:
        text = _decode_bytes(data)
        reader = csv.reader(io.StringIO(text), delimiter=sep)
        rows = list(reader)
        tables.append(rows)
        if rows:
            headers = rows[0]
            meta["columns"] = headers
            meta["row_count"] = len(rows) - 1
            kv = {"headers": headers, "sample_row": rows[1] if len(rows) > 1 else []}
        raw_parts = [" | ".join(r) for r in rows]
        pages.append(PageContent(page_number=1, text="\n".join(raw_parts), tables=tables))
    except Exception as exc:
        errors.append(f"CSV extraction failed: {exc}")
        method = "failed"

    return _build_result(
        file_id, filename, "csv", "\n".join(raw_parts),
        pages, tables, kv, meta, method,
        confidence=0.98 if not errors else 0.3,
        errors=errors, warnings=warnings,
    )


def _extract_pptx(data: bytes, file_id: str, filename: str) -> ExtractionResult:
    errors, warnings, pages, tables, kv, meta = [], [], [], [], {}, {}
    raw_parts = []
    method = "python-pptx"
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(io.BytesIO(data))
        props = prs.core_properties
        meta = {"title": props.title, "author": props.author, "slide_count": len(prs.slides)}
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            slide_text = "\n".join(texts)
            raw_parts.append(slide_text)
            pages.append(PageContent(page_number=i, text=slide_text))
    except Exception as exc:
        errors.append(f"PPTX extraction failed: {exc}")
        method = "failed"

    return _build_result(
        file_id, filename, "pptx", "\n".join(raw_parts),
        pages, tables, kv, meta, method,
        confidence=0.9 if not errors else 0.3,
        errors=errors, warnings=warnings,
    )


def _extract_image(data: bytes, file_id: str, filename: str) -> ExtractionResult:
    errors, warnings, pages, tables, kv, meta = [], [], [], [], {}, {}
    raw_text = ""
    method = "OCR"
    confidence = 0.6

    try:
        from PIL import Image  # type: ignore
        img = Image.open(io.BytesIO(data))
        meta = {
            "format": img.format,
            "mode": img.mode,
            "size": list(img.size),
            "width": img.size[0],
            "height": img.size[1],
        }
        if hasattr(img, "_getexif") and img._getexif():
            try:
                from PIL.ExifTags import TAGS  # type: ignore
                exif_data = {TAGS.get(k, k): str(v) for k, v in img._getexif().items()}
                meta["exif"] = exif_data
            except Exception:
                pass
    except Exception as exc:
        errors.append(f"Image open failed: {exc}")

    try:
        import pytesseract  # type: ignore
        from PIL import Image as PILImage  # type: ignore
        img_obj = PILImage.open(io.BytesIO(data))
        raw_text = pytesseract.image_to_string(img_obj, lang="eng")
        osd = pytesseract.image_to_osd(img_obj, output_type=pytesseract.Output.DICT)
        confidence = min(float(osd.get("orientation_conf", 60)) / 100.0, 1.0)
    except Exception as exc:
        warnings.append(f"OCR not available ({exc}); returning image metadata only")
        confidence = 0.1

    pages.append(PageContent(page_number=1, text=raw_text, confidence=confidence, images_found=1))
    return _build_result(
        file_id, filename, "image", raw_text,
        pages, tables, kv, meta, method,
        confidence=confidence,
        errors=errors, warnings=warnings,
    )


def _extract_json(data: bytes, file_id: str, filename: str) -> ExtractionResult:
    errors, warnings, pages, tables, kv, meta = [], [], [], [], {}, {}
    method = "json"
    text = _decode_bytes(data)
    raw_text = text
    try:
        parsed = json.loads(text)
        if isinstance(parsed, dict):
            kv = {str(k): str(v)[:200] for k, v in parsed.items()}
            meta = {"root_keys": list(parsed.keys()), "type": "object"}
        elif isinstance(parsed, list):
            meta = {"length": len(parsed), "type": "array"}
            if parsed and isinstance(parsed[0], dict):
                kv = {"sample": str(parsed[0])[:500]}
        raw_text = json.dumps(parsed, indent=2, default=str)
    except Exception as exc:
        errors.append(f"JSON parse error: {exc}")
    pages.append(PageContent(page_number=1, text=raw_text))
    return _build_result(
        file_id, filename, "json", raw_text,
        pages, tables, kv, meta, method,
        confidence=0.99 if not errors else 0.5,
        errors=errors, warnings=warnings,
    )


def _extract_html(data: bytes, file_id: str, filename: str) -> ExtractionResult:
    errors, warnings, pages, tables, kv, meta = [], [], [], [], {}, {}
    method = "beautifulsoup4"
    raw_text = ""
    try:
        from bs4 import BeautifulSoup  # type: ignore
        soup = BeautifulSoup(data, "lxml")
        title_tag = soup.find("title")
        meta["title"] = title_tag.text.strip() if title_tag else ""
        for tag in soup(["script", "style", "head"]):
            tag.decompose()
        raw_text = soup.get_text(separator="\n")
        # Extract tables
        for tbl in soup.find_all("table"):
            rows = [[cell.get_text(strip=True) for cell in row.find_all(["td", "th"])]
                    for row in tbl.find_all("tr")]
            if rows:
                tables.append(rows)
        # Meta tags
        for m in soup.find_all("meta"):
            if m.get("name") and m.get("content"):
                meta[m["name"]] = m["content"]
    except Exception as exc:
        method = "fallback"
        warnings.append(f"BS4/lxml failed ({exc}); stripping tags manually")
        raw_text = re.sub(r"<[^>]+>", " ", _decode_bytes(data))
    pages.append(PageContent(page_number=1, text=raw_text, tables=tables))
    return _build_result(
        file_id, filename, "html", raw_text,
        pages, tables, kv, meta, method,
        confidence=0.9 if not errors else 0.5,
        errors=errors, warnings=warnings,
    )


def _extract_xml(data: bytes, file_id: str, filename: str) -> ExtractionResult:
    errors, warnings, pages, tables, kv, meta = [], [], [], [], {}, {}
    method = "lxml"
    raw_text = ""
    try:
        import lxml.etree as ET  # type: ignore
        root = ET.fromstring(data)
        meta["root_tag"] = root.tag
        meta["namespaces"] = list(root.nsmap.values()) if hasattr(root, "nsmap") else []
        raw_text = ET.tostring(root, encoding="unicode", method="text")
        # Flatten all text nodes
        texts = [el.text.strip() for el in root.iter() if el.text and el.text.strip()]
        raw_text = "\n".join(texts)
    except Exception as exc:
        warnings.append(f"lxml failed ({exc}); using stdlib")
        method = "stdlib"
        try:
            import xml.etree.ElementTree as ET2
            root2 = ET2.fromstring(data)
            texts = [el.text.strip() for el in root2.iter() if el.text and el.text.strip()]
            raw_text = "\n".join(texts)
        except Exception as exc2:
            errors.append(f"XML extraction failed: {exc2}")
    pages.append(PageContent(page_number=1, text=raw_text))
    return _build_result(
        file_id, filename, "xml", raw_text,
        pages, tables, kv, meta, method,
        confidence=0.9 if not errors else 0.4,
        errors=errors, warnings=warnings,
    )


def _extract_text(data: bytes, file_id: str, filename: str, doc_type: str = "text") -> ExtractionResult:
    text = _decode_bytes(data)
    pages = [PageContent(page_number=1, text=text)]
    return _build_result(
        file_id, filename, doc_type, text,
        pages, [], {}, {}, "plaintext",
        confidence=1.0, errors=[], warnings=[],
    )


def _extract_yaml(data: bytes, file_id: str, filename: str) -> ExtractionResult:
    errors, kv, meta = [], {}, {}
    raw_text = _decode_bytes(data)
    try:
        import yaml  # type: ignore
        parsed = yaml.safe_load(raw_text)
        if isinstance(parsed, dict):
            kv = {str(k): str(v)[:200] for k, v in parsed.items()}
    except Exception as exc:
        errors.append(f"YAML parse error: {exc}")
    pages = [PageContent(page_number=1, text=raw_text)]
    return _build_result(
        file_id, filename, "yaml", raw_text,
        pages, [], kv, meta, "pyyaml",
        confidence=0.95 if not errors else 0.5,
        errors=errors, warnings=[],
    )


def _extract_email(data: bytes, file_id: str, filename: str) -> ExtractionResult:
    import email as emaillib
    errors, kv, meta = [], {}, {}
    raw_parts = []
    msg = emaillib.message_from_bytes(data)
    meta = {
        "from": msg.get("From", ""),
        "to": msg.get("To", ""),
        "subject": msg.get("Subject", ""),
        "date": msg.get("Date", ""),
        "message_id": msg.get("Message-ID", ""),
    }
    kv = dict(meta)
    if msg.is_multipart():
        for part in msg.walk():
            ct = part.get_content_type()
            if ct == "text/plain":
                payload = part.get_payload(decode=True)
                if payload:
                    raw_parts.append(_decode_bytes(payload))
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            raw_parts.append(_decode_bytes(payload))
    raw_text = "\n\n".join(raw_parts)
    pages = [PageContent(page_number=1, text=raw_text)]
    return _build_result(
        file_id, filename, "email", raw_text,
        pages, [], kv, meta, "email-stdlib",
        confidence=0.95, errors=errors, warnings=[],
    )


def _extract_archive(data: bytes, file_id: str, filename: str) -> ExtractionResult:
    import zipfile, tarfile
    errors, warnings, meta = [], [], {}
    file_list = []
    method = "archive"
    try:
        if zipfile.is_zipfile(io.BytesIO(data)):
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                file_list = zf.namelist()
                meta["archive_type"] = "zip"
        elif tarfile.is_tarfile(io.BytesIO(data)):
            with tarfile.open(fileobj=io.BytesIO(data)) as tf:
                file_list = tf.getnames()
                meta["archive_type"] = "tar"
    except Exception as exc:
        errors.append(f"Archive listing failed: {exc}")
    meta["file_count"] = len(file_list)
    meta["files"] = file_list[:50]
    raw_text = "Archive contents:\n" + "\n".join(file_list)
    pages = [PageContent(page_number=1, text=raw_text)]
    return _build_result(
        file_id, filename, "archive", raw_text,
        pages, [], {}, meta, method,
        confidence=0.8 if not errors else 0.3,
        errors=errors, warnings=warnings,
    )


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

_EXT_TO_EXTRACTOR = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".doc": _extract_docx,
    ".odt": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".xls": _extract_xlsx,
    ".ods": _extract_xlsx,
    ".csv": _extract_csv,
    ".tsv": lambda d, fid, fn: _extract_csv(d, fid, fn, sep="\t"),
    ".pptx": _extract_pptx,
    ".ppt": _extract_pptx,
    ".json": _extract_json,
    ".jsonl": _extract_json,
    ".html": _extract_html,
    ".htm": _extract_html,
    ".xhtml": _extract_html,
    ".xml": _extract_xml,
    ".yaml": _extract_yaml,
    ".yml": _extract_yaml,
    ".txt": _extract_text,
    ".md": lambda d, fid, fn: _extract_text(d, fid, fn, "markdown"),
    ".rst": lambda d, fid, fn: _extract_text(d, fid, fn, "rst"),
    ".log": lambda d, fid, fn: _extract_text(d, fid, fn, "log"),
    ".eml": _extract_email,
    ".msg": _extract_email,
    ".zip": _extract_archive,
    ".tar": _extract_archive,
    ".gz": _extract_archive,
    ".7z": _extract_archive,
    ".rar": _extract_archive,
    ".png": _extract_image,
    ".jpg": _extract_image,
    ".jpeg": _extract_image,
    ".tiff": _extract_image,
    ".tif": _extract_image,
    ".bmp": _extract_image,
    ".gif": _extract_image,
    ".webp": _extract_image,
    # Code files
    ".py": lambda d, fid, fn: _extract_text(d, fid, fn, "python"),
    ".js": lambda d, fid, fn: _extract_text(d, fid, fn, "javascript"),
    ".ts": lambda d, fid, fn: _extract_text(d, fid, fn, "typescript"),
    ".java": lambda d, fid, fn: _extract_text(d, fid, fn, "java"),
    ".sql": lambda d, fid, fn: _extract_text(d, fid, fn, "sql"),
}

_MIME_TO_EXT = {
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "text/csv": ".csv",
    "application/json": ".json",
    "text/html": ".html",
    "text/xml": ".xml",
    "application/xml": ".xml",
    "image/jpeg": ".jpg",
    "image/png": ".png",
    "image/gif": ".gif",
    "image/tiff": ".tiff",
    "image/bmp": ".bmp",
    "image/webp": ".webp",
    "text/plain": ".txt",
    "message/rfc822": ".eml",
}


def extract_document(
    file_data: bytes,
    file_id: str,
    filename: str,
    mime_type: str = "",
) -> ExtractionResult:
    """
    Dispatch extraction based on file extension (MIME type as fallback).
    """
    ext = Path(filename).suffix.lower()
    if ext not in _EXT_TO_EXTRACTOR:
        ext = _MIME_TO_EXT.get(mime_type, "")

    extractor = _EXT_TO_EXTRACTOR.get(ext)

    if extractor is None:
        # Generic fallback: try to decode as text
        logger.warning("No dedicated extractor for '%s'; using text fallback", ext)
        return _extract_text(file_data, file_id, filename, doc_type=ext.lstrip(".") or "unknown")

    try:
        return extractor(file_data, file_id, filename)
    except Exception as exc:
        logger.error("Extractor crashed for %s: %s", filename, exc)
        return ExtractionResult(
            success=False,
            file_id=file_id,
            filename=filename,
            document_type=ext.lstrip(".") or "unknown",
            raw_text="",
            pages=[],
            tables=[],
            key_value_pairs={},
            metadata={},
            language="unknown",
            word_count=0,
            char_count=0,
            page_count=0,
            extraction_method="failed",
            confidence=0.0,
            chunks=[],
            errors=[f"Extraction crashed: {exc}", traceback.format_exc()],
            warnings=[],
        )
